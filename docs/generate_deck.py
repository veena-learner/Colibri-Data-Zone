"""
Generate Colibri Data Zone C-Level Deck as PPTX
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colors ──────────────────────────────────────────────────
NAVY      = RGBColor(0x0D, 0x1B, 0x2A)
NAVY2     = RGBColor(0x1A, 0x2E, 0x44)
NAVY3     = RGBColor(0x13, 0x34, 0x4F)
TEAL      = RGBColor(0x0A, 0x7E, 0x7E)
TEAL2     = RGBColor(0x0D, 0x9B, 0x9B)
TEAL3     = RGBColor(0xB2, 0xE0, 0xE0)
GOLD      = RGBColor(0xE8, 0xA8, 0x35)
GOLD2     = RGBColor(0xF5, 0xC0, 0x57)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE  = RGBColor(0xF4, 0xF7, 0xF9)
GRAY      = RGBColor(0x6B, 0x72, 0x80)
LGRAY     = RGBColor(0xE5, 0xE9, 0xEE)
MUTED     = RGBColor(0x8B, 0x9A, 0xAA)
GREEN     = RGBColor(0x16, 0xA3, 0x4A)
DANGER    = RGBColor(0xD9, 0x40, 0x40)
LBLUE     = RGBColor(0x22, 0x56, 0x8A)

# Slide size 16:9
SW = Inches(13.333)
SH = Inches(7.5)

# ── Helpers ─────────────────────────────────────────────────
def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill, line=None, alpha=None):
    shp = slide.shapes.add_shape(1, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    return shp


def txbox(slide, text, l, t, w, h, size, color,
          bold=False, italic=False, align=PP_ALIGN.LEFT,
          wrap=True, name='Calibri'):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = name
    return box


def multiline_txbox(slide, lines, l, t, w, h, size, color,
                    bold=False, align=PP_ALIGN.LEFT, line_space=1.2,
                    name='Calibri'):
    """lines = list of (text, bold, color, size) or plain strings"""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    from pptx.util import Pt as PPt
    from pptx.oxml.ns import qn
    from lxml import etree

    first = True
    for item in lines:
        if isinstance(item, str):
            txt, b, c, s = item, bold, color, size
        else:
            txt, b, c, s = item

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = txt
        run.font.size = PPt(s)
        run.font.color.rgb = c
        run.font.bold = b
        run.font.name = name
    return box


def header_block(slide, eyebrow, title, subtitle=None,
                 eyebrow_color=TEAL2, title_color=WHITE, sub_color=MUTED,
                 l=Inches(0.8), t=Inches(0.55), w=Inches(11.5)):
    """Renders eyebrow + title + optional subtitle"""
    txbox(slide, eyebrow, l, t, w, Inches(0.35),
          10, eyebrow_color, bold=True)
    txbox(slide, title, l, t + Inches(0.38), w, Inches(0.85),
          30, title_color, bold=True)
    if subtitle:
        txbox(slide, subtitle, l, t + Inches(1.25), w, Inches(0.6),
              13, sub_color)


def bottom_bar(slide, text, bg=NAVY2, fg=MUTED):
    rect(slide, 0, SH - Inches(0.55), SW, Inches(0.55), bg)
    txbox(slide, text, Inches(0.5), SH - Inches(0.5),
          SW - Inches(1), Inches(0.45), 10, fg,
          align=PP_ALIGN.CENTER)


def left_accent(slide, color=TEAL):
    rect(slide, 0, 0, Inches(0.22), SH, color)


def card(slide, l, t, w, h, bg, title, body,
         title_color=WHITE, body_color=MUTED,
         icon=None, title_size=13, body_size=11, border=None):
    rect(slide, l, t, w, h, bg, line=border)
    cur_t = t + Inches(0.18)
    if icon:
        txbox(slide, icon, l + Inches(0.18), cur_t, Inches(0.45), Inches(0.4),
              18, title_color)
        txbox(slide, title, l + Inches(0.62), cur_t + Inches(0.01),
              w - Inches(0.8), Inches(0.35), title_size, title_color, bold=True)
        cur_t += Inches(0.42)
    else:
        txbox(slide, title, l + Inches(0.18), cur_t,
              w - Inches(0.36), Inches(0.35), title_size, title_color, bold=True)
        cur_t += Inches(0.38)
    txbox(slide, body, l + Inches(0.18), cur_t,
          w - Inches(0.36), h - (cur_t - t) - Inches(0.15),
          body_size, body_color, wrap=True)


# ════════════════════════════════════════════════════════════
# BUILD PRESENTATION
# ════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH

# ────────────────────────────────────────────────────────────
# SLIDE 1 — Cover
# ────────────────────────────────────────────────────────────
s = new_slide(prs)
set_bg(s, NAVY)
left_accent(s)

# Diagonal decorative gradient shape
rect(s, Inches(8.5), 0, Inches(5), SH, NAVY2)
rect(s, Inches(10), 0, Inches(3.5), SH, NAVY3)

# Horizontal teal line
rect(s, Inches(0.22), Inches(3.8), Inches(7), Inches(0.06), TEAL)

# Confidential badge top-right
txbox(s, 'CONFIDENTIAL  |  C-LEVEL BRIEFING  |  MARCH 2026',
      Inches(6), Inches(0.28), Inches(7), Inches(0.35),
      9, MUTED, align=PP_ALIGN.RIGHT)

# Brand
box = s.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(9), Inches(1.2))
tf = box.text_frame
tf.word_wrap = False
p = tf.paragraphs[0]
for txt, col in [('Colibri ', WHITE), ('Data Zone', TEAL2)]:
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(56)
    r.font.bold = True
    r.font.name = 'Calibri'
    r.font.color.rgb = col

txbox(s, 'An Enterprise Data Governance Portal',
      Inches(0.6), Inches(2.85), Inches(9), Inches(0.45),
      18, TEAL3)
txbox(s, 'Built for Colibri Group',
      Inches(0.6), Inches(3.35), Inches(7), Inches(0.38),
      13, GOLD)
txbox(s, 'One place to discover, govern, and trace data from source systems\nto the warehouse — the foundation for trustworthy analytics and AI.',
      Inches(0.6), Inches(3.85), Inches(8.5), Inches(0.75),
      13, MUTED)

# Meta row
for i, (label, val) in enumerate([('Organization', 'Colibri Group'),
                                    ('Presentation Date', 'March 2026'),
                                    ('Distribution', 'Confidential')]):
    x = Inches(0.6) + i * Inches(3.0)
    txbox(s, label, x, Inches(5.2), Inches(2.6), Inches(0.3), 9, TEAL2, bold=True)
    txbox(s, val,   x, Inches(5.52), Inches(2.6), Inches(0.3), 12, WHITE)

bottom_bar(s, 'Colibri Group   |   March 2026   |   Confidential', NAVY2, MUTED)

# ────────────────────────────────────────────────────────────
# SLIDE 2 — The Problem
# ────────────────────────────────────────────────────────────
s = new_slide(prs)
set_bg(s, OFFWHITE)
left_accent(s)

header_block(s,
    '01 — THE PROBLEM',
    "Why Data Governance Can't Wait",
    "As Colibri Group scales, fragmented metadata and uncontrolled data access have become operational and\ncompliance risks that directly slow down analytics, AI adoption, and decision-making.",
    eyebrow_color=TEAL, title_color=NAVY, sub_color=GRAY)

# 5 problem cards (2 cols)
problems = [
    ('🗂', 'Fragmented Metadata',
     'Data definitions, ownership, and business rules live in spreadsheets, wikis, and tribal knowledge — causing conflicting reports and eroded trust.'),
    ('🔐', 'Inconsistent Access',
     'No single view of who can access what. Audit requests require manual reconstruction across multiple systems with no auditability.'),
    ('🔍', 'No Lineage Visibility',
     'When a report number looks wrong, root-cause analysis is manual and time-consuming. There is no automated path from output back to source.'),
    ('📦', 'Duplicated Assets',
     'Multiple copies of datasets exist across teams without clear provenance, increasing costs and risk of decisions made on stale data.'),
    ('⚠️', 'Compliance Exposure',
     'Without centralized governance, demonstrating regulatory compliance requires manual, error-prone exercises consuming significant legal and stewardship resources.'),
]
card_w = Inches(4.9)
card_h = Inches(1.35)
gap    = Inches(0.22)
start_l = Inches(0.5)
start_t = Inches(2.05)

for i, (icon, title, body) in enumerate(problems):
    col = i % 3
    row = i // 3
    l = start_l + col * (card_w + gap)
    t = start_t + row * (card_h + gap)
    if i == 4:  # 5th card — center in row 2
        l = Inches(0.5) + 1 * (card_w + gap)
    card(s, l, t, card_w, card_h, WHITE, title, body,
         title_color=NAVY, body_color=GRAY, icon=icon,
         title_size=12, body_size=10.5, border=LGRAY)

# Quote
txbox(s, '"Without proper governance, even the most advanced analytics tools will produce unreliable, inconsistent, or unauthorized outputs."',
      Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.55),
      11, GRAY, italic=True, align=PP_ALIGN.CENTER)

bottom_bar(s, 'Colibri Data Zone   |   C-Level Briefing   |   March 2026', LGRAY, GRAY)

# ────────────────────────────────────────────────────────────
# SLIDE 3 — The Solution
# ────────────────────────────────────────────────────────────
s = new_slide(prs)
set_bg(s, NAVY3)
left_accent(s, TEAL2)

header_block(s,
    '02 — THE SOLUTION',
    'Introducing Colibri Data Zone',
    'A purpose-built data governance portal — inspired by AWS DataZone architecture, tailored to the Colibri Group\noperating environment — providing the governance foundation for trusted analytics and AI.',
    eyebrow_color=TEAL2, title_color=WHITE, sub_color=TEAL3)

solutions = [
    ('🔎', 'Self-Service Discovery',
     'A searchable, governed catalog of all data assets across business domains. Business users find trusted data in minutes, not hours.'),
    ('🏛', 'Federated Governance',
     'Domain-based model maps to business units. Each domain owns its assets, policies, and glossary — with enterprise-wide visibility.'),
    ('🔗', 'End-to-End Lineage',
     'Automated tracking via dbt manifest parsing, from OLTP source systems through Fivetran and dbt transformations to Redshift.'),
    ('📜', 'Business Glossary',
     'A shared vocabulary of terms, definitions, and synonyms ensuring consistent interpretation across all users — including AI agents.'),
    ('🔐', 'Governed Access',
     'Role-based access controls with publish-subscribe workflows. Every access decision is documented, auditable, and revocable.'),
    ('🤖', 'AI-Ready Foundation',
     'The semantic layer, access controls, and governed catalog that Conversational AI agents need to return accurate, authorized results.'),
]
card_w = Inches(3.9)
card_h = Inches(1.45)
gap    = Inches(0.22)
start_l = Inches(0.4)
start_t = Inches(2.05)

for i, (icon, title, body) in enumerate(solutions):
    col = i % 3
    row = i // 3
    l = start_l + col * (card_w + gap)
    t = start_t + row * (card_h + gap)
    card(s, l, t, card_w, card_h, NAVY2, title, body,
         title_color=TEAL2, body_color=TEAL3, icon=icon,
         title_size=12, body_size=10.5)

bottom_bar(s, 'Colibri Data Zone   |   C-Level Briefing   |   March 2026', NAVY, TEAL3)

# ────────────────────────────────────────────────────────────
# SLIDE 4 — Architecture
# ────────────────────────────────────────────────────────────
s = new_slide(prs)
set_bg(s, WHITE)
left_accent(s)

txbox(s, '03 — ARCHITECTURE', Inches(0.5), Inches(0.3), Inches(11), Inches(0.35),
      10, TEAL, bold=True)
txbox(s, 'Built for Our Stack. Deployed in Our Environment.',
      Inches(0.5), Inches(0.68), Inches(11), Inches(0.6),
      26, NAVY, bold=True)

# 4-layer architecture stack
layers = [
    (TEAL,  WHITE,  'PRESENTATION LAYER',
     'React 18 + TypeScript browser portal — search, catalog, lineage visualization, governance admin'),
    (TEAL2, WHITE,  'SERVICE LAYER',
     'Node.js/Express REST API — catalog ops, access management, glossary, lineage computation, workflow orchestration'),
    (NAVY2, TEAL3,  'INTEGRATION LAYER',
     'Connectors to Redshift, S3, dbt manifest/catalog, Fivetran, identity providers — metadata sync + lineage graphs'),
    (NAVY,  TEAL3,  'DATA LAYER',
     'AWS DynamoDB (single-table design) — metadata store, policy engine, audit log, lineage graph persistence'),
]
lw = Inches(7.2)
lh = Inches(0.78)
lx = Inches(0.4)
ly = Inches(1.42)
gap = Inches(0.1)

for i, (bg, fg, lbl, desc) in enumerate(layers):
    rect(s, lx, ly + i*(lh+gap), lw, lh, bg)
    txbox(s, lbl,  lx + Inches(0.25), ly + i*(lh+gap) + Inches(0.1),
          Inches(2.8), Inches(0.35), 11, fg, bold=True)
    txbox(s, desc, lx + Inches(0.25), ly + i*(lh+gap) + Inches(0.38),
          lw - Inches(0.5), Inches(0.35), 11, fg)

# Tech stack right panel
rect(s, Inches(7.9), Inches(1.35), Inches(5.2), Inches(5.3), OFFWHITE, line=LGRAY)
txbox(s, 'Technology Stack', Inches(8.1), Inches(1.5), Inches(4.8), Inches(0.4),
      14, NAVY, bold=True)

tech = [
    ('⚛️', 'React 18 + Vite + TypeScript', 'TanStack Query, React Flow, Tailwind CSS'),
    ('🟩', 'Node.js / Express / TypeScript', 'JWT auth, role-based access middleware'),
    ('🗄', 'AWS DynamoDB', 'Single-table design, PITR, SSE, pay-per-request'),
    ('☁️', 'AWS ECS Fargate + ALB', 'Containerized; Docker Compose for local dev'),
    ('🔧', 'dbt Integration', 'Parses manifest.json + catalog.json for lineage DAG'),
]
for i, (icon, name, detail) in enumerate(tech):
    ty = Inches(2.0) + i * Inches(0.88)
    txbox(s, icon, Inches(8.1), ty, Inches(0.45), Inches(0.4), 18, NAVY)
    txbox(s, name,   Inches(8.6), ty,              Inches(4.2), Inches(0.3), 12, NAVY, bold=True)
    txbox(s, detail, Inches(8.6), ty + Inches(0.3), Inches(4.2), Inches(0.28), 10, GRAY)

bottom_bar(s, 'Colibri Data Zone   |   C-Level Briefing   |   March 2026', LGRAY, GRAY)

# ────────────────────────────────────────────────────────────
# SLIDE 5 — Platform Capabilities
# ────────────────────────────────────────────────────────────
s = new_slide(prs)
set_bg(s, NAVY)
left_accent(s, GOLD)

header_block(s,
    '04 — PLATFORM CAPABILITIES',
    'What Colibri Data Zone Does Today',
    subtitle=None,
    eyebrow_color=GOLD, title_color=WHITE)

capabilities = [
    ('📚', 'Data Catalog',
     'Register, search, browse, and govern data assets across all domains. Grid/list views, sensitivity tagging, schema documentation, and bulk Excel upload.'),
    ('🏢', 'Domain Management',
     'Hierarchical business domain structure with owners, stewards, and tags. Finance, Marketing, Engineering, LMS, Sales, and more.'),
    ('📖', 'Business Glossary',
     'Shared definitions, synonyms, and related terms linked directly to catalog assets. The semantic layer for AI and analytics alike.'),
    ('🕸', 'Data Lineage Engine',
     'Interactive dbt lineage graph from manifest.json. Manual asset-to-asset lineage. Upstream/downstream impact analysis with configurable depth traversal.'),
    ('🗂', 'Source Table Catalog',
     'PM persona for OLTP-to-Redshift ingestion lifecycle. Tracks ColibriLMS, Salesforce, NetSuite, HubSpot, Zendesk, and ColibriPlatform tables.'),
    ('🧬', 'Ontology Layer',
     'Column-level descriptions, enhanced definitions, and ontology class for all dbt model columns. Bulk upload via CSV/Excel with upsert semantics.'),
]
card_w = Inches(3.9)
card_h = Inches(1.5)
gap    = Inches(0.2)
start_l = Inches(0.4)
start_t = Inches(1.75)

for i, (icon, title, body) in enumerate(capabilities):
    col = i % 3
    row = i // 3
    l = start_l + col * (card_w + gap)
    t = start_t + row * (card_h + gap)
    # "Live" badge
    bg_card = NAVY2
    card(s, l, t, card_w, card_h, bg_card, title, body,
         title_color=TEAL2, body_color=TEAL3, icon=icon,
         title_size=12, body_size=10.5)
    # Live badge
    rect(s, l + card_w - Inches(0.75), t + Inches(0.12), Inches(0.6), Inches(0.25), GREEN)
    txbox(s, 'Live', l + card_w - Inches(0.72), t + Inches(0.13),
          Inches(0.55), Inches(0.22), 9, WHITE, bold=True, align=PP_ALIGN.CENTER)

bottom_bar(s, 'Colibri Data Zone   |   C-Level Briefing   |   March 2026', NAVY2, TEAL3)

# ────────────────────────────────────────────────────────────
# SLIDE 6 — Business Benefits
# ────────────────────────────────────────────────────────────
s = new_slide(prs)
set_bg(s, OFFWHITE)
left_accent(s)

header_block(s,
    '05 — BUSINESS BENEFITS',
    'Measurable Value Across Three Dimensions',
    subtitle=None,
    eyebrow_color=TEAL, title_color=NAVY)

columns = [
    ('⚡', 'Operational Efficiency', TEAL, [
        'Eliminate ad hoc data access requests and manual provisioning',
        'Reduce time to find the right data from hours to minutes',
        'Streamline onboarding for new analysts and business users',
        'Bulk upload workflows replace one-by-one data entry',
        'Self-service analytics without engineering tickets',
    ]),
    ('🏅', 'Trust & Quality', GOLD, [
        'Single source of truth for all data definitions and business terms',
        'Sensitivity classification for every asset — Public to Restricted',
        'Lineage visibility enables rapid root-cause analysis',
        'Freshness and ownership indicators build data confidence',
        'Quality scores surface trusted versus unreliable datasets',
    ]),
    ('🚀', 'Strategic Enablement', NAVY2, [
        'Foundation for Conversational AI over governed, trusted data',
        'Federated model aligned to business unit autonomy',
        'Audit-ready compliance without manual reconstruction',
        'Positions Colibri Group for scalable, compliant data growth',
        'No vendor lock-in; full deployment control',
    ]),
]
col_w = Inches(4.0)
col_h = Inches(3.8)
start_l = Inches(0.4)
start_t = Inches(1.85)
gap = Inches(0.27)

for i, (icon, title, accent, bullets) in enumerate(columns):
    l = start_l + i * (col_w + gap)
    rect(s, l, start_t, col_w, col_h, WHITE, line=LGRAY)
    # Accent top bar
    rect(s, l, start_t, col_w, Inches(0.12), accent)
    txbox(s, icon,  l + Inches(0.2), start_t + Inches(0.22), Inches(0.5), Inches(0.4), 20, NAVY)
    txbox(s, title, l + Inches(0.7), start_t + Inches(0.25), col_w - Inches(0.85), Inches(0.38), 14, NAVY, bold=True)
    for j, bul in enumerate(bullets):
        txbox(s, f'• {bul}',
              l + Inches(0.2), start_t + Inches(0.75) + j * Inches(0.56),
              col_w - Inches(0.35), Inches(0.52), 10.5, GRAY)

# Stats row
stats = [('0', 'Per-User Licensing Costs'), ('Weeks', 'vs. Months for Commercial Tools'),
         ('6', 'User Roles — Viewer to Admin'), ('Full', 'Deployment & Data Control')]
stat_w = Inches(3.0)
for i, (num, label) in enumerate(stats):
    l = Inches(0.4) + i * Inches(3.23)
    t = Inches(5.8)
    txbox(s, num,   l, t,              stat_w, Inches(0.5), 26, TEAL, bold=True)
    txbox(s, label, l, t + Inches(0.48), stat_w, Inches(0.35), 10, GRAY)

bottom_bar(s, 'Colibri Data Zone   |   C-Level Briefing   |   March 2026', LGRAY, GRAY)

# ────────────────────────────────────────────────────────────
# SLIDE 7 — Use Cases
# ────────────────────────────────────────────────────────────
s = new_slide(prs)
set_bg(s, NAVY3)
left_accent(s, GOLD)

header_block(s,
    '06 — USE CASES',
    'Five Critical Business Scenarios Enabled',
    subtitle=None,
    eyebrow_color=GOLD, title_color=WHITE)

use_cases = [
    ('🔎', '1. Self-Service Data Discovery',
     'Business users search the portal to find trusted, governed datasets without knowing database schemas. Quality scores, ownership, and freshness guide selection.'),
    ('🤝', '2. Governed Data Sharing Across BUs',
     'Data producers publish curated data products. Other business units discover and subscribe via documented, auditable, revocable workflows — replacing ad hoc exports and grants.'),
    ('🤖', '3. Conversational AI Governance Layer',
     'The glossary ensures AI interprets terms correctly; the catalog defines approved datasets; access policies ensure AI respects permissions — making AI analytics trustworthy.'),
    ('🔭', '4. Impact Analysis & Change Management',
     'Before schema changes or model deprecation, lineage reveals full downstream impact — preventing accidental breakage of dashboards, reports, and AI agents.'),
    ('📋', '5. Compliance & Audit Readiness',
     'A complete, queryable audit trail of who accessed what data, when, and for what purpose. Regulatory evidence is produced from the portal — not reconstructed manually.'),
]
card_w = Inches(5.5)
card_h = Inches(1.45)
gap    = Inches(0.22)
start_l = [Inches(0.4), Inches(6.15)]
start_t = Inches(1.75)

for i, (icon, title, body) in enumerate(use_cases):
    col = i % 2
    row = i // 2
    if i == 4:  # 5th card — centered
        l = Inches(3.65)
        t = start_t + 2 * (card_h + gap)
    else:
        l = start_l[col]
        t = start_t + row * (card_h + gap)
    card(s, l, t, card_w, card_h, NAVY2, title, body,
         title_color=GOLD, body_color=TEAL3, icon=icon,
         title_size=12, body_size=10.5)

bottom_bar(s, 'Colibri Data Zone   |   C-Level Briefing   |   March 2026', NAVY, TEAL3)

# ────────────────────────────────────────────────────────────
# SLIDE 8 — Market Comparison
# ────────────────────────────────────────────────────────────
s = new_slide(prs)
set_bg(s, WHITE)
left_accent(s)

txbox(s, '07 — MARKET COMPARISON', Inches(0.5), Inches(0.3), Inches(11), Inches(0.35),
      10, TEAL, bold=True)
txbox(s, 'How We Compare to Market Alternatives',
      Inches(0.5), Inches(0.68), Inches(11), Inches(0.6), 26, NAVY, bold=True)
txbox(s, 'Evaluated: AWS DataZone, Collibra, Alation, Atlan',
      Inches(0.5), Inches(1.3), Inches(11), Inches(0.35), 12, GRAY)

# Table
from pptx.util import Inches as I
rows = [
    ['Capability',              'Colibri Data Zone', 'AWS DataZone', 'Collibra',   'Alation',    'Atlan'],
    ['Data Catalog',            '✓',                '✓',            '✓',          '✓',          '✓'],
    ['Business Glossary',       '✓',                '✓',            '✓',          '✓',          '✓'],
    ['Data Lineage (dbt)',       '✓ Native',         'Partial',      'Connector',  'Limited',    '✓'],
    ['Access Request Workflows','✓',                '✓',            '✓',          'Limited',    '✓'],
    ['Custom-Built for Stack',  '✓',                '✗',            '✗',          '✗',          '✗'],
    ['Redshift Integration',    'Deep Native',      'Deep',         'Connector',  'Connector',  'Connector'],
    ['Per-User Licensing',      'None ✓',           'Usage-Based',  'Enterprise', 'Enterprise', 'Enterprise'],
    ['Deployment Control',      'Full ✓',           'AWS Only',     'SaaS',       'SaaS',       'SaaS Only'],
    ['Implementation',          'Weeks ✓',          'Weeks',        '6-12 Months','Months',     'Months'],
    ['Ongoing Cost',            'Low ✓',            'Moderate',     'High',       'High',       'Moderate'],
]

tbl_left = Inches(0.35)
tbl_top  = Inches(1.72)
tbl_w    = Inches(12.6)
tbl_h    = Inches(5.2)
col_widths = [Inches(2.6), Inches(2.1), Inches(1.85), Inches(1.75), Inches(1.75), Inches(1.75)]

table = s.shapes.add_table(len(rows), 6, tbl_left, tbl_top, tbl_w, tbl_h).table

for ci, cw in enumerate(col_widths):
    table.columns[ci].width = cw

for ri, row_data in enumerate(rows):
    for ci, cell_text in enumerate(row_data):
        cell = table.cell(ri, ci)
        cell.text = cell_text
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(10.5)
        run.font.name = 'Calibri'
        # Header row
        if ri == 0:
            run.font.bold = True
            run.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
        elif ci == 1:  # Colibri Data Zone column highlight
            run.font.bold = True
            run.font.color.rgb = NAVY
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xE0, 0xF4, 0xF4)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 == 0 else OFFWHITE
            run.font.color.rgb = NAVY

bottom_bar(s, 'Colibri Data Zone   |   C-Level Briefing   |   March 2026', LGRAY, GRAY)

# ────────────────────────────────────────────────────────────
# SLIDE 9 — Our Advantage
# ────────────────────────────────────────────────────────────
s = new_slide(prs)
set_bg(s, NAVY)
left_accent(s, GOLD)

header_block(s,
    '08 — OUR ADVANTAGE',
    'Why Colibri Data Zone Is the Right Choice',
    subtitle=None,
    eyebrow_color=GOLD, title_color=WHITE)

advantages = [
    ('🎯', 'Tailored to Our Stack',
     'Built specifically for our Redshift-centric architecture and dbt transformation layer. No connector complexity, no impedance mismatch with generic platforms.'),
    ('💰', 'Zero Per-User Licensing',
     'Critical for data democratization. Every leader and business user can access governed data without costs scaling linearly with organizational adoption.'),
    ('🏗', 'Proven Architectural Patterns',
     "Modeled on AWS DataZone's battle-tested design (domains, data products, publish-subscribe) — without the constraints or lock-in of a managed AWS service."),
    ('🎛', 'Full Customization Control',
     'We adapt governance policies, organizational structure, and UX to our exact needs. No off-the-shelf product allows this degree of tailoring.'),
    ('⚡', 'Rapid Time to Value',
     'Already built and deployed — iterative enhancements take weeks, not the 6-12 months required by commercial governance platforms like Collibra.'),
    ('🤖', 'Direct AI Strategy Synergy',
     'The governance portal directly feeds the Conversational AI strategy by providing semantic context, access controls, and the trusted data catalog AI agents need.'),
]
card_w = Inches(3.9)
card_h = Inches(1.5)
gap    = Inches(0.2)
start_l = Inches(0.4)
start_t = Inches(1.75)

for i, (icon, title, body) in enumerate(advantages):
    col = i % 3
    row = i // 3
    l = start_l + col * (card_w + gap)
    t = start_t + row * (card_h + gap)
    card(s, l, t, card_w, card_h, NAVY2, title, body,
         title_color=GOLD, body_color=TEAL3, icon=icon,
         title_size=12, body_size=10.5)

bottom_bar(s, 'Colibri Data Zone   |   C-Level Briefing   |   March 2026', NAVY2, TEAL3)

# ────────────────────────────────────────────────────────────
# SLIDE 10 — Strategic Context / AI
# ────────────────────────────────────────────────────────────
s = new_slide(prs)
set_bg(s, OFFWHITE)
left_accent(s)

txbox(s, '09 — STRATEGIC CONTEXT', Inches(0.5), Inches(0.3), Inches(11), Inches(0.35),
      10, TEAL, bold=True)
txbox(s, 'Governance Enables Trustworthy Conversational AI',
      Inches(0.5), Inches(0.68), Inches(11), Inches(0.6), 26, NAVY, bold=True)
txbox(s, 'Colibri Data Zone and Conversational AI are complementary investments. The governance portal ensures that the data AI queries is trusted, well-defined, and properly permissioned. Without governance, conversational AI risks returning inaccurate or unauthorized results.',
      Inches(0.5), Inches(1.35), Inches(12.3), Inches(0.65), 12, GRAY)

# Pipeline flow
pipeline = [
    ('🗄', 'Source Systems', 'ColibriLMS, Salesforce,\nNetSuite, HubSpot'),
    ('🔧', 'dbt + Redshift', 'Transformed, governed\nwarehouse models'),
    ('🏛', 'Colibri Data Zone', 'Catalog · Glossary\nLineage · Policies'),
    ('🤖', 'AI Analytics Agent', 'Conversational queries\nover governed data'),
    ('📊', 'Business Leaders', 'Trusted, authorized,\naccurate insights'),
]
pipe_w = Inches(2.0)
pipe_h = Inches(1.35)
pipe_t = Inches(2.1)
pipe_gap = Inches(0.45)
total_pipe = len(pipeline) * pipe_w + (len(pipeline)-1) * pipe_gap
pipe_l_start = (float(SW) - total_pipe) / 2

for i, (icon, title, detail) in enumerate(pipeline):
    l = pipe_l_start + i * (pipe_w + pipe_gap)
    bg = NAVY2 if i == 2 else NAVY
    rect(s, l, pipe_t, pipe_w, pipe_h, bg)
    txbox(s, icon,   l + Inches(0.15), pipe_t + Inches(0.12), pipe_w - Inches(0.3), Inches(0.4), 20, TEAL2)
    txbox(s, title,  l + Inches(0.15), pipe_t + Inches(0.48), pipe_w - Inches(0.3), Inches(0.35), 11, WHITE, bold=True)
    txbox(s, detail, l + Inches(0.15), pipe_t + Inches(0.82), pipe_w - Inches(0.3), Inches(0.45), 9, TEAL3)
    # Arrow
    if i < len(pipeline) - 1:
        arr_l = l + pipe_w + Inches(0.05)
        txbox(s, '→', arr_l, pipe_t + Inches(0.5), Inches(0.35), Inches(0.4), 18, TEAL)

# 3 pillars
pillars = [
    ('📖', 'Semantic Accuracy',
     'Business glossary ensures AI interprets "Revenue", "Churn Rate", and "Customer" consistently with organizational definitions — not its training data defaults.'),
    ('🔐', 'Authorized Access',
     'Access control policies in the governance portal govern which datasets the AI agent can query — ensuring sensitive data is never surfaced to unauthorized users.'),
    ('✅', 'Trusted Data Foundation',
     'Only catalog-registered, steward-approved datasets are surfaced to the AI — eliminating the risk of AI answers based on stale, duplicate, or ungoverned data.'),
]
pillar_w = Inches(3.9)
pillar_h = Inches(1.3)
pillar_t = Inches(3.65)
pillar_gap = Inches(0.27)
pillar_l_start = Inches(0.4)

for i, (icon, title, body) in enumerate(pillars):
    l = pillar_l_start + i * (pillar_w + pillar_gap)
    card(s, l, pillar_t, pillar_w, pillar_h, WHITE, title, body,
         title_color=NAVY, body_color=GRAY, icon=icon,
         title_size=12, body_size=10.5, border=LGRAY)

bottom_bar(s, 'Colibri Data Zone   |   C-Level Briefing   |   March 2026', LGRAY, GRAY)

# ────────────────────────────────────────────────────────────
# SLIDE 11 — Implementation Roadmap
# ────────────────────────────────────────────────────────────
s = new_slide(prs)
set_bg(s, WHITE)
left_accent(s)

txbox(s, '10 — IMPLEMENTATION ROADMAP', Inches(0.5), Inches(0.3), Inches(11), Inches(0.35),
      10, TEAL, bold=True)
txbox(s, 'A Phased Path to Enterprise-Wide Governance',
      Inches(0.5), Inches(0.68), Inches(11), Inches(0.6), 26, NAVY, bold=True)
txbox(s, 'Four phases — from foundational catalog to AI-integrated, enterprise-scale governance',
      Inches(0.5), Inches(1.3), Inches(11), Inches(0.35), 12, GRAY)

phases = [
    (1, 'Phase 1', 'Foundation', TEAL, [
        'Onboard governance domains',
        'Catalog priority Redshift datasets',
        'Populate business glossary',
        'Assign data owners & stewards',
    ], 'Searchable governed catalog\nfor priority assets'),
    (2, 'Phase 2', 'Access Governance', TEAL2, [
        'Enable publish-subscribe workflows',
        'Integrate SSO identity provider',
        'Enforce role-based access',
        'Begin audit trail tracking',
    ], 'Governed, auditable\ndata sharing'),
    (3, 'Phase 3', 'AI Integration', GOLD, [
        'Connect catalog to AI agents',
        'AI respects access policies',
        'Glossary-aligned AI terms',
        'AI lineage audit capabilities',
    ], 'AI analytics on governed,\ntrusted data'),
    (4, 'Phase 4', 'Scale & Mature', NAVY2, [
        'Expand to all business units',
        'Activate full lineage tracking',
        'Data quality scoring',
        'Stewardship processes & SLAs',
    ], 'Enterprise-wide governance\n& compliance'),
]
phase_w = Inches(2.85)
phase_h = Inches(3.8)
phase_t = Inches(1.75)
phase_gap = Inches(0.33)
phase_l = Inches(0.45)

for i, (num, ph_name, ph_title, color, bullets, outcome) in enumerate(phases):
    l = phase_l + i * (phase_w + phase_gap)
    # Header
    rect(s, l, phase_t, phase_w, Inches(0.6), color)
    txbox(s, f'{num}  {ph_name}', l + Inches(0.15), phase_t + Inches(0.12),
          phase_w - Inches(0.3), Inches(0.38), 13, WHITE, bold=True)
    # Body
    rect(s, l, phase_t + Inches(0.6), phase_w, phase_h - Inches(0.6), OFFWHITE, line=LGRAY)
    txbox(s, ph_title, l + Inches(0.15), phase_t + Inches(0.7),
          phase_w - Inches(0.3), Inches(0.38), 13, NAVY, bold=True)
    for j, bul in enumerate(bullets):
        txbox(s, f'• {bul}',
              l + Inches(0.15), phase_t + Inches(1.15) + j * Inches(0.45),
              phase_w - Inches(0.3), Inches(0.42), 10, GRAY)
    # Outcome
    rect(s, l, phase_t + Inches(3.45), phase_w, Inches(0.65), color)
    txbox(s, outcome, l + Inches(0.12), phase_t + Inches(3.5),
          phase_w - Inches(0.24), Inches(0.55), 10, WHITE, bold=True)

# Current status
rect(s, Inches(0.35), Inches(5.7), Inches(12.5), Inches(0.62), RGBColor(0xE0, 0xF4, 0xF4), line=TEAL)
txbox(s, '⚡ Current Status:  Phase 1 is complete and operational. All core entities (Assets, Domains, Glossary, Ontology, Lineage) are implemented with full CRUD APIs, role-based access, dbt lineage, and bulk upload workflows. Ready for Phase 2.',
      Inches(0.55), Inches(5.75), Inches(12.2), Inches(0.52), 11, TEAL, bold=False)

bottom_bar(s, 'Colibri Data Zone   |   C-Level Briefing   |   March 2026', LGRAY, GRAY)

# ────────────────────────────────────────────────────────────
# SLIDE 12 — Next Steps
# ────────────────────────────────────────────────────────────
s = new_slide(prs)
set_bg(s, NAVY)
left_accent(s, GOLD)

header_block(s,
    '11 — RECOMMENDED NEXT STEPS',
    'What We Are Asking For',
    'Colibri Data Zone is operational and ready to scale. The following decisions and investments will accelerate\nour path to enterprise-wide data governance and AI-readiness.',
    eyebrow_color=GOLD, title_color=WHITE, sub_color=TEAL3)

actions = [
    ('✅', 'Endorse the Platform',
     'Executive endorsement of Colibri Data Zone as the organization\'s official data governance standard — signaling organizational commitment and driving adoption across business units.'),
    ('👥', 'Appoint Domain Stewards',
     'Identify and formally appoint Data Owners and Data Stewards within each business domain. Governance only works when ownership is clear and accountable.'),
    ('📋', 'Prioritize the Catalog Backfill',
     'Allocate stewardship time to catalog the first 50–100 priority Redshift datasets, ensuring the catalog reflects the most-used, most-impactful data assets across the organization.'),
    ('🤖', 'Align with AI Roadmap',
     'Formally connect the data governance roadmap to the Conversational AI strategy. Ensure Phase 3 (AI Integration) is scoped into the AI agent development timeline.'),
]
act_w = Inches(5.8)
act_h = Inches(1.3)
gap   = Inches(0.2)
for i, (icon, title, body) in enumerate(actions):
    col = i % 2
    row = i // 2
    l = Inches(0.4) + col * (act_w + Inches(0.35))
    t = Inches(2.1) + row * (act_h + gap)
    card(s, l, t, act_w, act_h, NAVY2, title, body,
         title_color=GOLD, body_color=TEAL3, icon=icon,
         title_size=12, body_size=10.5)

# Quote
rect(s, Inches(0.4), Inches(5.72), Inches(12.3), Inches(0.6), TEAL)
txbox(s, '"Governance is not a constraint on innovation — it is the foundation that makes innovation trustworthy."',
      Inches(0.6), Inches(5.78), Inches(12.0), Inches(0.5),
      13, WHITE, italic=True, align=PP_ALIGN.CENTER)

bottom_bar(s, 'Colibri Data Zone   |   C-Level Briefing   |   March 2026', NAVY2, TEAL3)

# ────────────────────────────────────────────────────────────
# SLIDE 13 — Closing
# ────────────────────────────────────────────────────────────
s = new_slide(prs)
set_bg(s, NAVY)
left_accent(s, TEAL2)

rect(s, Inches(8.5), 0, Inches(5), SH, NAVY2)
rect(s, Inches(10), 0, Inches(3.5), SH, NAVY3)
rect(s, Inches(0.22), Inches(3.8), Inches(7), Inches(0.06), TEAL)

box = s.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(9), Inches(1.2))
tf = box.text_frame
p = tf.paragraphs[0]
for txt, col in [('Colibri ', WHITE), ('Data Zone', TEAL2)]:
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(52)
    r.font.bold = True
    r.font.name = 'Calibri'
    r.font.color.rgb = col

txbox(s, 'One place to discover, govern, and trace data\nfrom source systems to the warehouse',
      Inches(0.6), Inches(2.95), Inches(8), Inches(0.7), 16, TEAL3)

# Stats
stats = [('7', 'Portal Sections'), ('6', 'Role Types'), ('0', 'Licensing Fees'), ('4', 'Roadmap Phases')]
for i, (num, label) in enumerate(stats):
    l = Inches(0.6) + i * Inches(3.1)
    txbox(s, num,   l, Inches(4.0), Inches(2.6), Inches(0.65), 40, GOLD, bold=True)
    txbox(s, label, l, Inches(4.7), Inches(2.6), Inches(0.35), 12, TEAL3)

bottom_bar(s, 'Colibri Group   |   March 2026   |   Confidential', NAVY2, MUTED)

# ════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════
out = '/Users/veena.anantharam/Project-Dev/Colibri-Data-zone/docs/Colibri_Data_Zone_CLevel_Deck.pptx'
prs.save(out)
print(f'Saved: {out}')
